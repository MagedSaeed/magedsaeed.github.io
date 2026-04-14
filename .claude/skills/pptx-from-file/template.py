"""
template.py — Atlas Navy template for the pptx-from-file skill.

Design philosophy (out of the box, no user tuning required):
  - Multi-tone palette: dark navy bg, cream ink, amber accent, teal secondary,
    muted blue-grey supporting text — richness from tonal variation, not a
    single accent.
  - Sans-serif typography with strong weight/size hierarchy.
  - Signature motifs repeated across layouts:
      * trio of colored marks (amber / teal / amber rectangles) as brand sig
      * amber eyebrow in small caps above every header
      * first-word underline on content headers
      * solid amber vertical-bar bullet markers (no glyph bullets)
      * teal-dash footer with presentation title + author on every content slide
      * "NN / TT" slide counter in amber top-right
      * ghosted huge numeral on section dividers (darker navy, not accent-filled)

API (compatible with prior skill usage):

    from template import Presentation, Theme

    pres = Presentation(theme=Theme.default(), subtitle="Subtitle line")
    pres.set_animation("points")
    pres.title_slide(title="...", eyebrow="A GUIDE", subtitle="...",
                     author="...", role="...", date="2026 · v1")
    pres.section_divider(number="01", title="...", subtitle="...")
    pres.content(title="...", bullets=[...], eyebrow="PART · CONCEPT",
                 callout_title="NOTE", callout_body="...")
    pres.two_column(title="...", bullets=[...], image_path="...", caption="...")
    pres.figure_full(image_path="...", caption="...")
    pres.quote(text="...", attribution="...", eyebrow="SECTION")
    pres.closing(text="Thank You.", subtitle="...", url="...")
    pres.save(out)

Legacy signatures (content without eyebrow, etc.) still work — all new params
are optional.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation as _PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# --------------------------------------------------------------------------- #
# Palette
# --------------------------------------------------------------------------- #

# Atlas Navy defaults (out-of-the-box design)
NAVY_BG    = "#0F1B2E"   # deep navy background
NAVY_GHOST = "#1B2A42"   # slightly lighter — ghosted numerals, tinted panels
CREAM_INK  = "#F5EFE3"   # warm cream, primary text
MUTED      = "#8B9BB0"   # blue-grey, supporting text
AMBER      = "#F5A524"   # primary accent — eyebrows, bullets, counter
TEAL       = "#3EC4B1"   # secondary accent — footer tick, secondary marks
PANEL_BG   = "#14223A"   # callout panel fill
RULE       = "#2A3A56"   # hairlines

# Alternate accents available for domain switches
ACCENTS = {
    "amber":  "#F5A524",
    "warm":   "#E26D4F",
    "forest": "#3EA47A",
    "indigo": "#7684D6",
    "plum":   "#B76CAB",
    "teal":   "#3EC4B1",
}

# Light-palette companion (flipped, same philosophy)
LIGHT_BG    = "#F7F3EB"
LIGHT_GHOST = "#ECE6D8"
LIGHT_INK   = "#16253D"
LIGHT_MUTED = "#6B7688"
LIGHT_PANEL = "#EEE7D6"
LIGHT_RULE  = "#D6CDBA"


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# --------------------------------------------------------------------------- #
# Theme
# --------------------------------------------------------------------------- #

@dataclass
class Theme:
    bg: str
    ghost: str
    ink: str
    muted: str
    accent: str
    accent2: str
    panel: str
    rule: str
    header_font: str = "Lato"
    body_font: str = "Lato"
    mono_font: str = "Consolas"
    is_dark: bool = True

    @classmethod
    def default(cls, accent: str = "amber", accent2: str = "teal") -> "Theme":
        """The out-of-the-box Atlas Navy theme."""
        return cls(
            bg=NAVY_BG, ghost=NAVY_GHOST, ink=CREAM_INK, muted=MUTED,
            accent=ACCENTS.get(accent, AMBER),
            accent2=ACCENTS.get(accent2, TEAL),
            panel=PANEL_BG, rule=RULE, is_dark=True,
        )

    @classmethod
    def dark(cls, accent: str = "amber", header_font: str = "Lato") -> "Theme":
        t = cls.default(accent=accent)
        t.header_font = header_font
        t.body_font = header_font
        return t

    @classmethod
    def light(cls, accent: str = "warm", header_font: str = "Lato") -> "Theme":
        return cls(
            bg=LIGHT_BG, ghost=LIGHT_GHOST, ink=LIGHT_INK, muted=LIGHT_MUTED,
            accent=ACCENTS.get(accent, "#C2410C"),
            accent2=ACCENTS.get("forest", "#3EA47A"),
            panel=LIGHT_PANEL, rule=LIGHT_RULE,
            header_font=header_font, body_font=header_font, is_dark=False,
        )

    # --- Preset variants — same philosophy, different expression ------------
    # Use these to avoid the "every presentation looks the same" problem.
    # Each preset is a deliberate pairing of palette + motif tuning.

    @classmethod
    def atlas_navy(cls) -> "Theme":
        """Deep navy + amber + teal. Tech / ML / research / product."""
        return cls.default(accent="amber", accent2="teal")

    @classmethod
    def forest(cls) -> "Theme":
        """Dark forest green + gold + terracotta. Nature / biology / eco."""
        return cls(
            bg="#0E2118", ghost="#17301F", ink="#F1ECD9", muted="#8FA294",
            accent="#D6A83C", accent2="#C2633B",
            panel="#152A1C", rule="#2A4030", is_dark=True,
        )

    @classmethod
    def plum(cls) -> "Theme":
        """Dark plum + coral + electric blue. Creative / design / pitch."""
        return cls(
            bg="#1C1024", ghost="#2A1B36", ink="#F3E8EF", muted="#A898B0",
            accent="#FF7A6A", accent2="#7AA7FF",
            panel="#241530", rule="#3A2444", is_dark=True,
        )

    @classmethod
    def slate(cls) -> "Theme":
        """Cool slate + cobalt + bronze. Corporate / formal / finance."""
        return cls(
            bg="#1A2230", ghost="#263142", ink="#EDEFF3", muted="#8791A3",
            accent="#5A7EFF", accent2="#CC8B4A",
            panel="#1F2A3B", rule="#2E3A4E", is_dark=True,
        )

    @classmethod
    def kraft(cls) -> "Theme":
        """Warm kraft paper + burnt orange + forest. Academic / humanities."""
        return cls(
            bg="#EFE6D2", ghost="#E4D9BD", ink="#2A1F15", muted="#6B5B45",
            accent="#B54A1F", accent2="#3C6E4F",
            panel="#E7DAB9", rule="#CFBF9B",
            header_font="Lato", body_font="Lato", is_dark=False,
        )

    @classmethod
    def sand(cls) -> "Theme":
        """Warm sand + deep navy + coral. Light variant, design-friendly."""
        return cls(
            bg="#F5EDDD", ghost="#EBE1CB", ink="#0F1B2E", muted="#6B7688",
            accent="#E15A50", accent2="#2D7A6E",
            panel="#EBE1C5", rule="#D8CCAF",
            header_font="Lato", body_font="Lato", is_dark=False,
        )


# --------------------------------------------------------------------------- #
# Layout math — 16:9, 13.333 x 7.5 in
# --------------------------------------------------------------------------- #

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)
GUTTER  = Inches(0.35)

BODY_LEFT  = MARGIN
BODY_RIGHT = SLIDE_W - MARGIN
BODY_WIDTH = BODY_RIGHT - BODY_LEFT


# --------------------------------------------------------------------------- #
# Animation XML helpers — byParagraph on-click reveal
# --------------------------------------------------------------------------- #

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _para_appear_timing_xml(shape_id: int, preset_id: int = 1, preset_class: str = "entr") -> str:
    return f"""
<p:timing xmlns:p="{_P_NS}" xmlns:a="{_A_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="{preset_id}" presetClass="{preset_class}"
                                     presetSubtype="0" fill="hold" nodeType="clickEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="1" fill="hold">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl><p:spTgt spid="{shape_id}"><p:txEl><p:pRg st="0" end="999"/></p:txEl></p:spTgt></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    <p:bldP spid="{shape_id}" grpId="0" build="byParagraph"/>
  </p:bldLst>
</p:timing>
""".strip()


_ANIMATION_PRESETS = {
    "points":     (1,  "entr"),
    "expressive": (10, "entr"),
    "none":       None,
}


def _attach_paragraph_animation(slide, shape, mode: str) -> None:
    if mode == "none" or mode not in _ANIMATION_PRESETS or _ANIMATION_PRESETS[mode] is None:
        return
    preset_id, preset_class = _ANIMATION_PRESETS[mode]
    timing_xml = _para_appear_timing_xml(shape.shape_id, preset_id, preset_class)
    timing_el = etree.fromstring(timing_xml)
    sld = slide._element
    existing = sld.find(qn("p:timing"))
    if existing is not None:
        sld.remove(existing)
    sld.append(timing_el)


# --------------------------------------------------------------------------- #
# Presentation wrapper
# --------------------------------------------------------------------------- #

class Presentation:
    """Atlas Navy presentation builder — rich multi-tone design by default."""

    def __init__(self, theme: Optional[Theme] = None,
                 subtitle: str = "", style: str = "rich"):
        if style not in ("rich", "minimal"):
            raise ValueError("style must be 'rich' or 'minimal'")
        self.theme = theme or Theme.default()
        self.style = style
        self.subtitle = subtitle               # used in footer
        self._footer_title = ""                # filled from title_slide
        self._footer_author = ""
        self.animation_mode = "points"

        self._pres = _PptxPresentation()
        self._pres.slide_width = SLIDE_W
        self._pres.slide_height = SLIDE_H
        self._blank_layout = self._pres.slide_layouts[6]

        self._slides_meta: list[dict] = []
        self._counter_refs: list[tuple] = []   # (textbox, current_index)

    # ---------- API ----------

    def set_animation(self, mode: str) -> None:
        if mode not in _ANIMATION_PRESETS:
            raise ValueError(f"mode must be one of {list(_ANIMATION_PRESETS)}")
        self.animation_mode = mode

    # ---------- title ----------

    def title_slide(self, title: str, eyebrow: str = "",
                    subtitle: str = "", author: str = "",
                    role: str = "", date: str = "") -> None:
        slide = self._new_slide(show_number=False, show_footer=False)
        self._footer_title = title
        self._footer_author = author

        # Trio of brand marks, top-left (amber square, teal bar, amber rect)
        self._draw_brand_trio(slide, top=Inches(0.55), left=MARGIN)

        # Amber eyebrow
        if eyebrow:
            self._text(slide, eyebrow.upper(),
                       left=MARGIN, top=Inches(2.1),
                       width=Inches(10), height=Inches(0.35),
                       font=self.theme.body_font, size=13,
                       color=self.theme.accent,
                       bold=True, letter_spacing=300)

        # Huge title (two-line friendly)
        self._text(slide, title,
                   left=MARGIN, top=Inches(2.55),
                   width=Inches(11), height=Inches(2.4),
                   font=self.theme.header_font, size=60,
                   color=self.theme.ink, bold=True, line_spacing=1.05)

        # Subtitle in muted blue-grey
        if subtitle:
            self._text(slide, subtitle,
                       left=MARGIN, top=Inches(5.1),
                       width=Inches(11), height=Inches(0.5),
                       font=self.theme.body_font, size=18,
                       color=self.theme.muted, line_spacing=1.3)

        # Author block: small amber bar + name + role
        author_top = Inches(6.1)
        if author:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, MARGIN, author_top + Inches(0.1),
                Inches(0.3), Inches(0.05)
            )
            bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(self.theme.accent)
            bar.line.fill.background()

            self._text(slide, author.upper(),
                       left=MARGIN + Inches(0.45), top=author_top,
                       width=Inches(6), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.ink, bold=True, letter_spacing=200)
            if role:
                self._text(slide, role,
                           left=MARGIN + Inches(0.45), top=author_top + Inches(0.3),
                           width=Inches(6), height=Inches(0.3),
                           font=self.theme.body_font, size=10,
                           color=self.theme.muted)

        # Date bottom-right
        if date:
            self._text(slide, date,
                       left=SLIDE_W - Inches(3), top=SLIDE_H - Inches(0.55),
                       width=Inches(2.5), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.muted, align=PP_ALIGN.RIGHT)

        self._slides_meta.append({"kind": "title", "title": title})

    # ---------- section divider ----------

    def section_divider(self, number: str, title: str,
                        subtitle: str = "", eyebrow: str = "PART") -> None:
        slide = self._new_slide(show_number=False, show_footer=False)

        # Giant ghosted numeral on the left (slightly lighter than bg)
        self._text(slide, number,
                   left=Inches(0.4), top=Inches(1.4),
                   width=Inches(6.5), height=Inches(5),
                   font=self.theme.header_font, size=380,
                   color=self.theme.ghost, bold=True, line_spacing=0.95)

        # Right side: amber bar + eyebrow
        bar_top = Inches(3.2)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(7.1), bar_top,
            Inches(0.5), Inches(0.07)
        )
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(self.theme.accent)
        bar.line.fill.background()

        self._text(slide, eyebrow.upper(),
                   left=Inches(7.75), top=bar_top - Inches(0.1),
                   width=Inches(4), height=Inches(0.35),
                   font=self.theme.body_font, size=14,
                   color=self.theme.accent, bold=True, letter_spacing=300)

        # Large bold title
        self._text(slide, title,
                   left=Inches(7.1), top=Inches(3.7),
                   width=Inches(5.8), height=Inches(2.0),
                   font=self.theme.header_font, size=50,
                   color=self.theme.ink, bold=True, line_spacing=1.05)

        # Subtitle in muted
        if subtitle:
            self._text(slide, subtitle,
                       left=Inches(7.1), top=Inches(5.6),
                       width=Inches(5.8), height=Inches(1.2),
                       font=self.theme.body_font, size=17,
                       color=self.theme.muted, line_spacing=1.3)

        self._slides_meta.append({"kind": "section", "title": title})

    # ---------- content ----------

    def content(self, title: str, bullets: list[str],
                eyebrow: str = "",
                callout_title: str = "", callout_body: str = "",
                animation_override: Optional[str] = None) -> None:
        if len(bullets) > 6:
            raise ValueError("max 6 bullets per slide; split into two slides")
        slide = self._new_slide()

        # Eyebrow
        if eyebrow:
            self._text(slide, eyebrow.upper(),
                       left=MARGIN, top=Inches(0.45),
                       width=Inches(8), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.accent,
                       bold=True, letter_spacing=300)

        # Header with first-word underline
        header_top = Inches(0.9)
        self._header_with_underline(slide, title,
                                    left=MARGIN, top=header_top,
                                    width=BODY_WIDTH, size=36)

        # Body area
        body_top = Inches(2.15)

        has_callout = bool(callout_title or callout_body)
        if has_callout:
            # Left bullets column (wider), right callout panel
            col_w = (BODY_WIDTH - GUTTER) * 0.64
            right_left = BODY_LEFT + col_w + GUTTER
            right_w = BODY_WIDTH - col_w - GUTTER

            body = self._bullets(slide, bullets,
                                 left=BODY_LEFT, top=body_top,
                                 width=col_w, height=Inches(4.3))
            _attach_paragraph_animation(slide, body,
                                        animation_override or self.animation_mode)
            self._draw_callout(slide, callout_title, callout_body,
                               left=right_left, top=body_top,
                               width=right_w, height=Inches(4.3))
        else:
            body = self._bullets(slide, bullets,
                                 left=BODY_LEFT, top=body_top,
                                 width=BODY_WIDTH, height=Inches(4.5))
            _attach_paragraph_animation(slide, body,
                                        animation_override or self.animation_mode)

        self._slides_meta.append({"kind": "content", "title": title})

    # ---------- two_column ----------

    def two_column(self, title: str, bullets: list[str],
                   image_path: Optional[str] = None,
                   right_text: Optional[str] = None, caption: str = "",
                   eyebrow: str = "",
                   animation_override: Optional[str] = None) -> None:
        if len(bullets) > 5:
            raise ValueError("max 5 bullets in two_column; split")
        slide = self._new_slide()

        if eyebrow:
            self._text(slide, eyebrow.upper(),
                       left=MARGIN, top=Inches(0.45),
                       width=Inches(8), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.accent,
                       bold=True, letter_spacing=300)

        self._header_with_underline(slide, title,
                                    left=MARGIN, top=Inches(0.9),
                                    width=BODY_WIDTH, size=34)

        body_top = Inches(2.15)
        col_w = (BODY_WIDTH - GUTTER) * 0.58
        right_left = BODY_LEFT + col_w + GUTTER
        right_w = BODY_WIDTH - col_w - GUTTER

        body = self._bullets(slide, bullets,
                             left=BODY_LEFT, top=body_top,
                             width=col_w, height=Inches(4.3))
        _attach_paragraph_animation(slide, body,
                                    animation_override or self.animation_mode)

        if image_path:
            self._image_with_caption(slide, image_path, caption,
                                     left=right_left, top=body_top,
                                     width=right_w, height=Inches(4.0))
        elif right_text:
            self._draw_callout(slide, "", right_text,
                               left=right_left, top=body_top,
                               width=right_w, height=Inches(4.3))

        self._slides_meta.append({"kind": "two_column", "title": title})

    # ---------- figure_full ----------

    def figure_full(self, image_path: str, caption: str = "",
                    title: str = "", eyebrow: str = "") -> None:
        slide = self._new_slide()
        if eyebrow:
            self._text(slide, eyebrow.upper(),
                       left=MARGIN, top=Inches(0.45),
                       width=Inches(8), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.accent,
                       bold=True, letter_spacing=300)
        if title:
            self._header_with_underline(slide, title,
                                        left=MARGIN, top=Inches(0.9),
                                        width=BODY_WIDTH, size=32)
            top = Inches(2.15); h = Inches(4.3)
        else:
            top = Inches(0.9); h = Inches(5.4)
        self._image_with_caption(slide, image_path, caption,
                                 left=MARGIN, top=top,
                                 width=BODY_WIDTH, height=h)
        self._slides_meta.append({"kind": "figure", "title": title or caption})

    # ---------- quote ----------

    def quote(self, text: str, attribution: str = "",
              eyebrow: str = "") -> None:
        slide = self._new_slide()

        if eyebrow:
            self._text(slide, eyebrow.upper(),
                       left=MARGIN, top=Inches(0.45),
                       width=Inches(8), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.accent,
                       bold=True, letter_spacing=300)

        # Amber short rule
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.1),
            Inches(0.6), Inches(0.06)
        )
        rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(self.theme.accent)
        rule.line.fill.background()

        # Amber ghost quote mark
        self._text(slide, "\u201C",
                   left=MARGIN, top=Inches(2.1),
                   width=Inches(1.5), height=Inches(2),
                   font=self.theme.header_font, size=160,
                   color=self.theme.accent, bold=True)

        # Huge bold cream quote body
        self._text(slide, text,
                   left=Inches(2.3), top=Inches(2.6),
                   width=Inches(10.5), height=Inches(3),
                   font=self.theme.header_font, size=40,
                   color=self.theme.ink, bold=True, line_spacing=1.2)

        if attribution:
            # Amber bar + attribution
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(5.85),
                Inches(0.35), Inches(0.05)
            )
            bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(self.theme.accent)
            bar.line.fill.background()
            self._text(slide, attribution,
                       left=Inches(2.8), top=Inches(5.75),
                       width=Inches(10), height=Inches(0.35),
                       font=self.theme.body_font, size=14,
                       color=self.theme.ink, bold=True, letter_spacing=100)

        self._slides_meta.append({"kind": "quote", "title": text[:40]})

    # ---------- closing ----------

    def closing(self, text: str = "Thank You.", subtitle: str = "",
                url: str = "") -> None:
        slide = self._new_slide(show_number=False, show_footer=False)

        # Amber square top-left
        mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.55),
            Inches(0.4), Inches(0.4)
        )
        mark.fill.solid(); mark.fill.fore_color.rgb = _rgb(self.theme.accent)
        mark.line.fill.background()

        # Huge cream "Thank You."
        self._text(slide, text,
                   left=MARGIN, top=Inches(2.2),
                   width=Inches(12), height=Inches(1.8),
                   font=self.theme.header_font, size=96,
                   color=self.theme.ink, bold=True, line_spacing=1.0)

        # Amber rule
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, Inches(4.15),
            Inches(0.9), Inches(0.05)
        )
        rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(self.theme.accent)
        rule.line.fill.background()

        if subtitle:
            self._text(slide, subtitle,
                       left=MARGIN, top=Inches(4.4),
                       width=Inches(12), height=Inches(0.5),
                       font=self.theme.body_font, size=18,
                       color=self.theme.muted, line_spacing=1.3)

        # Decorative diagonal wedge — a darker-navy triangle at bottom-right
        wedge = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(6.5), Inches(5.4),
            SLIDE_W - Inches(6.5), SLIDE_H - Inches(5.4)
        )
        wedge.fill.solid(); wedge.fill.fore_color.rgb = _rgb(self.theme.ghost)
        wedge.line.fill.background()
        # Flip horizontally so hypotenuse goes top-left -> bottom-right
        wedge.rotation = 0

        # URL in amber, bottom-left
        if url:
            self._text(slide, url,
                       left=MARGIN, top=SLIDE_H - Inches(0.7),
                       width=Inches(6), height=Inches(0.3),
                       font=self.theme.body_font, size=13,
                       color=self.theme.accent, bold=True)

        self._slides_meta.append({"kind": "closing", "title": text})

    # ---------- save ----------

    def save(self, path: str | Path) -> Path:
        # Fill in "NN / TT" counters now that we know total
        total = len(self._pres.slides)
        for (tb, idx) in self._counter_refs:
            tf = tb.text_frame
            p = tf.paragraphs[0]
            # rebuild runs: "NN / TT"
            # remove all existing runs first
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            r1 = p.add_run()
            r1.text = f"{idx:02d}"
            r1.font.name = self.theme.body_font
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = _rgb(self.theme.accent)
            r2 = p.add_run()
            r2.text = f" / {total:02d}"
            r2.font.name = self.theme.body_font
            r2.font.size = Pt(11)
            r2.font.bold = True
            r2.font.color.rgb = _rgb(self.theme.muted)

        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self._pres.save(str(path))
        return path

    # ---------- Internals ----------

    def _new_slide(self, show_number: bool = True, show_footer: bool = True):
        slide = self._pres.slides.add_slide(self._blank_layout)
        fill = slide.background.fill
        fill.solid(); fill.fore_color.rgb = _rgb(self.theme.bg)

        if show_number:
            n = len(self._pres.slides)
            tb = self._text(slide, f"{n:02d} / ??",
                            left=SLIDE_W - Inches(1.8),
                            top=Inches(0.45),
                            width=Inches(1.5), height=Inches(0.3),
                            font=self.theme.body_font, size=11,
                            color=self.theme.accent, bold=True,
                            align=PP_ALIGN.RIGHT)
            self._counter_refs.append((tb, n))

        if show_footer:
            self._draw_footer(slide)

        return slide

    def _draw_footer(self, slide):
        # teal dash
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, SLIDE_H - Inches(0.48),
            Inches(0.25), Inches(0.05)
        )
        tick.fill.solid(); tick.fill.fore_color.rgb = _rgb(self.theme.accent2)
        tick.line.fill.background()

        parts = [p for p in (self._footer_title or self.subtitle,
                             self._footer_author) if p]
        if parts:
            footer_text = "   ·   ".join(parts)
            self._text(slide, footer_text,
                       left=MARGIN + Inches(0.4),
                       top=SLIDE_H - Inches(0.55),
                       width=Inches(10), height=Inches(0.3),
                       font=self.theme.body_font, size=11,
                       color=self.theme.muted)

    def _draw_brand_trio(self, slide, top, left):
        """Trio of small marks: amber square, teal thin bar, amber bar."""
        # amber square
        s = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(0.28), Inches(0.28)
        )
        s.fill.solid(); s.fill.fore_color.rgb = _rgb(self.theme.accent)
        s.line.fill.background()
        # teal bar
        b1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left + Inches(0.4), top + Inches(0.1),
            Inches(0.08), Inches(0.2)
        )
        b1.fill.solid(); b1.fill.fore_color.rgb = _rgb(self.theme.accent2)
        b1.line.fill.background()
        # amber rectangle
        b2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left + Inches(0.6), top + Inches(0.08),
            Inches(0.35), Inches(0.2)
        )
        b2.fill.solid(); b2.fill.fore_color.rgb = _rgb(self.theme.accent)
        b2.line.fill.background()

    def _header_with_underline(self, slide, title: str, *,
                               left, top, width, size: int):
        """Header where the first word is underlined in accent color."""
        words = title.split(" ", 1)
        first = words[0]
        rest = " " + words[1] if len(words) > 1 else ""

        tb = slide.shapes.add_textbox(left, top, width, Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1

        r1 = p.add_run()
        r1.text = first
        r1.font.name = self.theme.header_font
        r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = _rgb(self.theme.ink)
        r1.font.underline = True
        # Set underline color to accent via XML
        rPr = r1._r.get_or_add_rPr()
        rPr.set("u", "sng")
        uFill = etree.SubElement(rPr, qn("a:uFill"))
        solidFill = etree.SubElement(uFill, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", self.theme.accent.lstrip("#"))

        if rest:
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = self.theme.header_font
            r2.font.size = Pt(size); r2.font.bold = True
            r2.font.color.rgb = _rgb(self.theme.ink)

        return tb

    def _draw_callout(self, slide, title: str, body: str, *,
                      left, top, width, height):
        """Tinted side panel with an eyebrow + body text (for callouts)."""
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        panel.fill.solid(); panel.fill.fore_color.rgb = _rgb(self.theme.panel)
        panel.line.fill.background()

        pad = Inches(0.3)
        inner_top = top + pad
        if title:
            self._text(slide, title.upper(),
                       left=left + pad, top=inner_top,
                       width=width - pad - pad, height=Inches(0.35),
                       font=self.theme.body_font, size=11,
                       color=self.theme.accent,
                       bold=True, letter_spacing=250)
            inner_top = inner_top + Inches(0.5)

        if body:
            self._text(slide, body,
                       left=left + pad, top=inner_top,
                       width=width - pad - pad, height=height - (inner_top - top) - pad,
                       font=self.theme.body_font, size=14,
                       color=self.theme.ink, line_spacing=1.4)

    def _text(self, slide, text: str, *, left, top, width, height,
              font: str, size: int, color: str,
              bold: bool = False, italic: bool = False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing: float = 1.2, letter_spacing: Optional[int] = None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = _rgb(color)
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
        return tb

    def _bullets(self, slide, bullets: list[str], *, left, top, width, height):
        """Vertical-bar marker bullets in accent, cream body text."""
        # draw bullet bars + text in a single textbox so animation can reveal
        # them per-paragraph; bar character is a full-block that renders as a
        # thick rectangle in most sans-serif fonts.
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        for i, txt in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.4
            p.space_after = Pt(10)
            # amber vertical-bar marker (full block char)
            r_mark = p.add_run()
            r_mark.text = "\u2588  "
            r_mark.font.name = self.theme.body_font
            r_mark.font.size = Pt(16); r_mark.font.bold = True
            r_mark.font.color.rgb = _rgb(self.theme.accent)
            # body text
            r = p.add_run()
            r.text = txt
            r.font.name = self.theme.body_font
            r.font.size = Pt(18)
            r.font.color.rgb = _rgb(self.theme.ink)
        return tb

    def _image_with_caption(self, slide, image_path: str, caption: str,
                            *, left, top, width, height):
        caption_h = Inches(0.45) if caption else Inches(0)
        img_h = height - caption_h
        pic = slide.shapes.add_picture(image_path, left, top,
                                       width=width, height=img_h)
        if caption:
            self._text(slide, caption,
                       left=left, top=top + img_h + Inches(0.1),
                       width=width, height=caption_h,
                       font=self.theme.body_font, size=12,
                       color=self.theme.muted, italic=True)
        return pic
